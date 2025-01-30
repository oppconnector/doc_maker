import os
from typing import Union
import PyPDF2
import docx
import csv
import openpyxl
import odf.opendocument
from odf.text import P
import striprtf
from bs4 import BeautifulSoup
import json
import xml.etree.ElementTree as ET
from ebooklib import epub
from pptx import Presentation
import pyexcel_ods

def read_doc(filepath: str) -> str:
    try:
        """
        Read a document file and return its content as a single string.
        
        :param filepath: The path to the document file
        :return: Content of the file as a string
        :raises ValueError: If the file type is not supported
        """
        _, file_extension = os.path.splitext(filepath)
        file_extension = file_extension.lower()


        if file_extension == '.pdf':
            print("doctype: PDF")
            with open(filepath, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ''
                for page in reader.pages:
                    text += page.extract_text()
            return text

        elif file_extension == '.docx':
            print("doctype: DOCX")
            doc = docx.Document(filepath)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            return '\n'.join(full_text)

        elif file_extension == '.txt':
            print("doctype: TXT")
            with open(filepath, 'r', encoding='utf-8') as file:
                return file.read()

        elif file_extension == '.csv':
            print("doctype: CSV")
            with open(filepath, 'r', newline='', encoding='utf-8') as file:
                reader = csv.reader(file)
                return '\n'.join([','.join(row) for row in reader])

        elif file_extension == '.xlsx':
            print("doctype: XLSX")
            workbook = openpyxl.load_workbook(filepath)
            sheet = workbook.active
            return '\n'.join(['\t'.join([str(cell.value or '') for cell in row]) for row in sheet.iter_rows()])

        elif file_extension == '.odt':
            print("doctype: ODT")
            doc = odf.opendocument.load(filepath)
            text = ''
            for p in doc.getElementsByType(P):
                text += p.getText() + '\n'
            return text

        elif file_extension == '.rtf':
            print("doctype: RTF")
            with open(filepath, 'r', encoding='utf-8') as file:
                rtf = file.read()
                return striprtf.strip_rtf(rtf)

        elif file_extension in ('.html', '.htm'):
            print("doctype: HTML")
            with open(filepath, 'r', encoding='utf-8') as file:
                soup = BeautifulSoup(file, 'html.parser')
                return soup.get_text()

        elif file_extension == '.json':
            print("doctype: JSON")
            with open(filepath, 'r', encoding='utf-8') as file:
                return json.dumps(json.load(file), indent=2)

        elif file_extension == '.xml':
            print("doctype: XML")
            tree = ET.parse(filepath)
            root = tree.getroot()
            return ET.tostring(root, encoding='unicode', method='text')

        elif file_extension == '.epub':
            print("doctype: EPUB")
            book = epub.read_epub(filepath)
            text = ''
            for doc in book.get_items_of_type(epub.ITEM_DOCUMENT):
                text += doc.get_body_content().decode('utf-8') + '\n'
            return text

        elif file_extension == '.pptx':
            print("doctype: PPTX")
            prs = Presentation(filepath)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
            return text

        elif file_extension == '.odp':
            print("doctype: ODP")
            doc = odf.opendocument.load(filepath)
            text = ''
            for slide in doc.getElementsByType('presentation:page'):
                for p in slide.getElementsByType('presentation:text'):
                    text += p.getText() + '\n'
            return text

        elif file_extension == '.ods':
            print("doctype: ODS")
            book = pyexcel_ods.get_data(filepath)
            text = ''
            for sheet_name, data in book.items():
                text += f"{sheet_name}:\n"
                for row in data:
                    text += '\t'.join(map(str, row)) + '\n'
            return text

        else:
            print(f"Could not read: {filepath}")
            return None
        
    except Exception as e:
        # Code to handle the exception
        print(f"An error occurred with the following file: {filepath}")
        print(f"Error: {e}")
        return None

# Example usage:
# text = read_doc('path/to/your/file.pdf')
# print(text)